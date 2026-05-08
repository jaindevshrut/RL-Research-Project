"""
Generates presentation/RL-Research-Project.pptx from the paper artefacts.

Produces an 8-slide deck:
    1. Title
    2. Motivation
    3. Method 1 — Dueling Double DQN with Bootstrap Ensemble
    4. Method 2 — QR-DQN (comparison)
    5. The 8-D Appraisal Vector
    6. Results — Master Conclusion Table + comparison bars
    7. Decorrelation evidence
    8. Conclusion and future work

Usage:  python presentation/make_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "paper" / "figures"
OUT = ROOT / "presentation" / "RL-Research-Project.pptx"

PRIMARY = RGBColor(0x1F, 0x77, 0xB4)
ACCENT  = RGBColor(0xFF, 0x7F, 0x0E)
DARK    = RGBColor(0x20, 0x20, 0x20)
GREY    = RGBColor(0x55, 0x55, 0x55)


def set_text(tf, text, size=18, bold=False, color=DARK):
    p = tf.paragraphs[0]
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                         Inches(12.3), Inches(0.9))
    set_text(title_box.text_frame, title, size=32, bold=True, color=PRIMARY)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(1.05),
                                     Inches(12.3), Inches(0.5))
        set_text(s.text_frame, subtitle, size=16, color=GREY)


def add_bullets(slide, items, left=0.5, top=1.7, width=8.0, height=5.5,
                font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return
    kwargs = {}
    if width is not None: kwargs["width"] = Inches(width)
    if height is not None: kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0),
                                   Inches(12.3), Inches(0.4))
    set_text(box.text_frame, text, size=10, color=GREY)


def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # -------------------------------------------------------------- Slide 1
    s = prs.slides.add_slide(blank)
    title_box = s.shapes.add_textbox(Inches(0.6), Inches(2.0),
                                     Inches(12.0), Inches(2.0))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("An Appraisal-Augmented Dueling Double DQN with a "
              "Bootstrap Ensemble for Cognitively Grounded Affective "
              "Modelling in Reinforcement Learning Agents")
    for r in p.runs:
        r.font.size = Pt(30); r.font.bold = True
        r.font.color.rgb = PRIMARY
    sub = s.shapes.add_textbox(Inches(0.6), Inches(4.4),
                               Inches(12.0), Inches(2.0))
    set_text(sub.text_frame,
             "Dev Shrut Jain  |  Sweta Rana  |  Krish Rathod",
             size=20, color=DARK)
    sub2 = s.shapes.add_textbox(Inches(0.6), Inches(4.9),
                                Inches(12.0), Inches(1.0))
    set_text(sub2.text_frame,
             "Department of Artificial Intelligence, SVNIT Surat",
             size=16, color=GREY)
    sub3 = s.shapes.add_textbox(Inches(0.6), Inches(5.4),
                                Inches(12.0), Inches(1.0))
    set_text(sub3.text_frame,
             "B.Tech. Project  |  Year 3",
             size=14, color=GREY)

    # -------------------------------------------------------------- Slide 2
    s = prs.slides.add_slide(blank)
    add_title(s, "Motivation",
              "Why couple Scherer's CPM with reinforcement learning?")
    add_bullets(s, [
        "Emotion is not a layer painted on top of cognition — it shares "
        "machinery with attention, action selection, and learning.",
        "Scherer's Component Process Model (CPM) decomposes emotion into "
        "a hierarchy of cognitive checks. Most of those checks already "
        "read quantities an RL agent maintains: TD error, Q-spread, "
        "transition counts.",
        "Zhang, Broekens, Jokinen (arXiv:2309.06367, 2023) operationalise "
        "four such checks — suddenness, goal relevance, conduciveness, "
        "power — on a tabular Q-learning agent.",
        "Limitation: three of those four checks ride the same TD signal. "
        "The 4-D vector cannot, by construction, distinguish emotions "
        "that differ on something other than TD magnitude (e.g. boredom "
        "vs neutral, hope vs joy).",
        "Our question: can the appraisal vector be enriched with channels "
        "whose informativeness is empirically verifiable, not asserted?",
    ], font_size=16)
    add_footer(s, "B.Tech. Project | SVNIT Surat | 2026")

    # -------------------------------------------------------------- Slide 3
    s = prs.slides.add_slide(blank)
    add_title(s, "Method 1 — Dueling Double DQN with Bootstrap Ensemble",
              "The headline backbone (chosen configuration)")
    add_bullets(s, [
        "Dueling head: Q(s,a) = V(s) + (A(s,a) − mean_a A). Exposes "
        "V(s) and the centred advantage as named outputs.",
        "Bootstrap ensemble: K = 3 independently-initialised value heads "
        "share a feature trunk. Cross-head dispersion σ_K(s) is a clean "
        "epistemic-uncertainty signal.",
        "Double-DQN target removes max-operator bias from the TD error.",
        "These three architectural pieces feed three appraisal channels "
        "directly: V → anticipation, advantage spread → power, "
        "σ_K → predictability.",
        "Trained for 60,000 frames on a 7×7 key-and-lava grid world; "
        "the appraisal vector is a read-out, the policy is unaffected.",
    ], font_size=16, width=7.5)
    add_image(s, FIGS / "training_curve.png", left=8.2, top=1.7, width=4.7)
    add_footer(s, "src/networks/dueling_dqn.py | src/agent.py | src/train.py")

    # -------------------------------------------------------------- Slide 4
    s = prs.slides.add_slide(blank)
    add_title(s, "Method 2 — Quantile-Regression DQN",
              "A controlled comparison (Dabney et al., AAAI 2018)")
    add_bullets(s, [
        "Replaces the dueling V-head and the K-ensemble with a single "
        "QR-DQN head producing N = 51 quantile atoms per action.",
        "Trained with the asymmetric quantile Huber loss "
        "(eq. 19 of the paper).",
        "Same seed, same env, same 60k-frame budget, same eight-dim "
        "appraisal vector — only the value head and loss are swapped.",
        "Three input substitutions: Q = mean over atoms; "
        "V(s) = max_a Q(s,a); predictability uses the within-distribution "
        "spread σ_Z of the greedy action's atoms.",
        "Why we ran it: to test whether Method 1's gains are specific "
        "to its representation. Result: yes, they are. Method 2 trails "
        "Method 1 on every appraisal-derived metric.",
    ], font_size=16)
    add_footer(s, "src/networks/qr_dqn.py | src/agent_qrdqn.py | src/train_qrdqn.py")

    # -------------------------------------------------------------- Slide 5
    s = prs.slides.add_slide(blank)
    add_title(s, "The 8-Dimensional Appraisal Vector",
              "Four CPM checks reproduced + four orthogonal extensions")
    add_bullets(s, [
        "x1 suddenness — 1 − P(s' | s, a)  [transition counts]",
        "x2 goal relevance — |TD| / running max",
        "x3 conduciveness — tanh(TD / scale)",
        "x4 power — centred Q-range / max|Q|",
        "x5 predictability★ — 1 − σ_K / running max  [ensemble dispersion]",
        "x6 anticipation★ — tanh(V(s) / scale)  [absolute value level]",
        "x7 urgency★ — t / T_max  [within-episode time index]",
        "x8 familiarity★ — log(N(s)+1) / log(N_max+1)  [state-visit count]",
        "★ = newly proposed; each computed from a statistic that is "
        "algebraically disjoint from the original four.",
    ], font_size=15, width=7.5)
    add_image(s, FIGS / "correlation_heatmap.png", left=8.2, top=1.7,
              width=4.7)
    add_footer(s, "src/appraisal/extractor.py")

    # -------------------------------------------------------------- Slide 6
    s = prs.slides.add_slide(blank)
    add_title(s, "Results — Master Conclusion Table",
              "60,000 frames, identical seed, three configurations")
    # Manual table
    rows = [
        ("Exp.",  "Configuration", "Backbone", "Acc.",  "R²",    "RMSE"),
        ("Exp 1", "Baseline-4D",   "Method 1", "0.648", "0.366", "0.345"),
        ("Exp 2", "Extended-8D",   "Method 1", "0.903", "0.757", "0.214"),
        ("Exp 3", "Extended-8D",   "Method 2", "0.868", "0.694", "0.239"),
    ]
    table = s.shapes.add_table(len(rows), len(rows[0]),
                               Inches(0.5), Inches(1.7),
                               Inches(7.6), Inches(2.6)).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.bold = (i == 0) or (i == 2)  # bold headers + Exp 2
                    r.font.color.rgb = PRIMARY if (i == 2) else DARK
    add_image(s, FIGS / "comparison_bars.png", left=8.2, top=1.7, width=4.7)
    note = s.shapes.add_textbox(Inches(0.5), Inches(4.5),
                                Inches(7.6), Inches(2.4))
    tf = note.text_frame; tf.word_wrap = True
    for line in [
        "• Extended-8D (Method 1) leads on every metric — the chosen configuration.",
        "• Method 2 (QR-DQN) is competitive but trails by 3.5 pp accuracy, "
        "0.063 R², 0.025 RMSE.",
        "• Effective rank of the appraisal covariance roughly doubles "
        "from 1.58/4 to 3.18/8.",
        "• Greedy task return identical between Baseline-4D and Method 1 "
        "(seed=0); appraisal layer is a read-out, not a policy modifier.",
    ]:
        p = tf.paragraphs[0] if line == "• Extended-8D (Method 1) leads on every metric — the chosen configuration." else tf.add_paragraph()
        p.text = line
        for r in p.runs:
            r.font.size = Pt(13); r.font.color.rgb = DARK
    add_footer(s, "runs/rmse_r2_summary.json | analysis/rmse_r2.py")

    # -------------------------------------------------------------- Slide 7
    s = prs.slides.add_slide(blank)
    add_title(s, "Decorrelation Evidence",
              "The new dimensions carry genuinely independent information")
    add_image(s, FIGS / "eigenvalue_spectrum.png", left=0.4, top=1.7,
              width=6.3)
    add_image(s, FIGS / "perm_importance.png",     left=6.9, top=1.7,
              width=6.0)
    note = s.shapes.add_textbox(Inches(0.4), Inches(5.4),
                                Inches(12.5), Inches(1.6))
    tf = note.text_frame; tf.word_wrap = True
    for i, line in enumerate([
        "• Eigenvalue spectrum (left): Baseline-4D collapses onto its first "
        "eigenvector (erank = 1.58/4); Extended-8D spreads mass over a "
        "longer tail (erank = 3.18/8).",
        "• Permutation importance (right): four of the five most "
        "discriminative dimensions in Extended-8D are the newly-introduced "
        "ones (familiarity, power, anticipation, predictability).",
        "• Per-dimension VIF: every channel except anticipation sits below 5; "
        "anticipation reaches 7.4 — reported transparently as a "
        "world-driven coupling, not a modelling oversight.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for r in p.runs:
            r.font.size = Pt(13); r.font.color.rgb = DARK
    add_footer(s, "analysis/correlation_analysis.py | analysis/explainability.py")

    # -------------------------------------------------------------- Slide 8
    s = prs.slides.add_slide(blank)
    add_title(s, "Conclusion and Future Scope",
              "What we showed, what we did not, what comes next")
    add_bullets(s, [
        "Replacing the tabular Q-learning backbone with a Dueling Double "
        "DQN + bootstrap ensemble lets us define four new appraisal "
        "channels, each from a statistic algebraically disjoint from "
        "the original four.",
        "Extended-8D (Method 1): accuracy 64.8% → 90.3%, R² 0.366 → 0.757, "
        "RMSE 0.345 → 0.214. Effective rank doubles. Greedy task return "
        "unchanged.",
        "QR-DQN (Method 2) is competitive but underperforms — this "
        "confirms the dueling V-head and the ensemble are doing real "
        "work rather than being incidental.",
        "Limitations: emotion labels are task-event proxies, not human "
        "vignette ratings; predictability is uninformative very early "
        "in training; tested on one grid world only.",
        "Future work: port to the vignette protocol of [Zhang et al. 2023]; "
        "joint train classifier and value network; close the loop by "
        "feeding appraisal back as intrinsic reward; replace count "
        "tables with hash-based pseudo-counts for high-dim envs.",
    ], font_size=15)
    add_footer(s,
               "Code & paper: github.com/jaindevshrut/RL-Research-Project")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    build()
